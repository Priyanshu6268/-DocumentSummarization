import os
import mimetypes
import ollama
from PIL import Image
from pptx import Presentation
import pandas as pd
from docx import Document as DocxDocument

from llama_index.core import VectorStoreIndex, StorageContext, load_index_from_storage
from llama_index.core.schema import Document
from llama_index.core.node_parser import SentenceSplitter
from llama_index.readers.file import PyMuPDFReader

from langchain_ollama.llms import OllamaLLM as LangchainOllama
from langchain_ollama import OllamaEmbeddings
from llama_index.llms.langchain import LangChainLLM
from llama_index.embeddings.langchain import LangchainEmbedding

import extract_msg
from llama_index.core import Settings

OLLAMA_LLM_MODEL = "llama3.2"
OLLAMA_EMBED_MODEL = "mxbai-embed-large"
OLLAMA_VISION_MODEL = "llava"

MSG_FILE_PATH = r"RedHat Partner Portal Login.msg"
PDF_FILE_PATH = r"Digital Certificate Generation.pdf"
IMAGE_FILE_PATH = r"Problem.jpg"
PPTX_FILE_PATH = r"58188_2024083009573563_1305058465_PPTX.pptx"
# TABULAR_FILE_PATH = r"clustered_alarms_5min_expanded_new11111.csv"
TABULAR_FILE_PATH = r"personality_dataset.xlsx"
DOCX_FILE_PATH = r"annexure-II-proforma-for-submission-of-nominations-for-indiaai-fellowship.docx"

STORAGE_DIR_MSG = "./storage_msg"
STORAGE_DIR_PDF = "./storage_pdf"
STORAGE_DIR_IMG = "./storage_img"
STORAGE_DIR_PPTX = "./storage_pptx"
STORAGE_DIR_TABULAR = "./storage_tabular"
STORAGE_DIR_DOCX = "./storage_docx"

def parse_outlook_msg(file_path):
    try:
        msg = extract_msg.Message(file_path)
        msg_message = (
            f"Subject: {msg.subject}\n"
            f"From: {msg.sender}\n"
            f"Date: {msg.date}\n\n"
            f"{msg.body}"
        )
        return [Document(text=msg_message)]
    except Exception as e:
        print(f"Error parsing MSG file {file_path}: {e}")
        return None

def parse_pdf(file_path):
    try:
        loader = PyMuPDFReader()
        documents = loader.load(file_path=file_path)
        return documents
    except Exception as e:
        print(f"Error parsing PDF file {file_path}: {e}")
        return None

def parse_image(file_path, vision_model=OLLAMA_VISION_MODEL):
    try:
        print(f"Sending image to Ollama Vision model: {vision_model} for {file_path}")
        response = ollama.chat(
            model=vision_model,
            messages=[{
                'role': 'user',
                'content': 'What is in this image? Summarize it.',
                'images': [file_path]
            }]
        )
        summary = response.get("message", {}).get("content", "")
        if not summary:
            print(f"Warning: No summary content returned for image {file_path}.")
        return [Document(text=summary)]
    except Exception as e:
        print(f"Error summarizing image {file_path}: {e}")
        return None

def parse_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            if slide_text:
                full_text.append(f"Slide {i+1}:\n" + "\n".join(slide_text))

        if not full_text:
            print(f"Warning: No text found in PPTX file {file_path}.")
            return None

        return [Document(text="\n\n".join(full_text))]
    except Exception as e:
        print(f"Error parsing PPTX file {file_path}: {e}")
        return None

def parse_tabular_data(file_path):
    try:
        df = None
        if file_path.lower().endswith('.csv'):
            df = pd.read_csv(file_path)
        elif file_path.lower().endswith(('.xls', '.xlsx')):
            df = pd.read_excel(file_path)
        else:
            print(f"Unsupported tabular file format: {file_path}")
            return None

        if df is None:
            return None

        insights = f"Data overview:\n"
        insights += f"Shape: {df.shape[0]} rows, {df.shape[1]} columns\n"
        insights += f"Columns: {', '.join(df.columns)}\n\n"

        numerical_cols = df.select_dtypes(include=['number']).columns
        if not numerical_cols.empty:
            insights += "Numerical column statistics:\n"
            insights += df[numerical_cols].describe().to_string() + "\n\n"

        categorical_cols = df.select_dtypes(include=['object', 'category']).columns
        if not categorical_cols.empty:
            insights += "Categorical column insights (top 5 unique values/counts):\n"
            for col in categorical_cols:
                insights += f"- {col}:\n"
                insights += df[col].value_counts().head(5).to_string() + "\n"
            insights += "\n"

        insights += "First 5 rows of data:\n"
        insights += df.head().to_string()

        return [Document(text=insights)]
    except Exception as e:
        print(f"Error parsing tabular file {file_path}: {e}")
        return None

def parse_docx(file_path):
    try:
        doc = DocxDocument(file_path)
        full_text = []
        for paragraph in doc.paragraphs:
            if paragraph.text:
                full_text.append(paragraph.text)
        
        if not full_text:
            print(f"Warning: No text found in DOCX file {file_path}.")
            return None
            
        return [Document(text="\n\n".join(full_text))]
    except Exception as e:
        print(f"Error parsing DOCX file {file_path}: {e}")
        return None

def process_document(file_path, storage_dir, document_type, is_image=False, is_tabular=False, is_word=False):
    print(f"\n--- Processing {document_type} file: {file_path} ---")
    os.makedirs(storage_dir, exist_ok=True)

    documents = None
    llm_model_to_use = OLLAMA_LLM_MODEL
    query_text = "Summarize the key messages or topics from the document."

    if document_type == "Outlook MSG":
        documents = parse_outlook_msg(file_path)
    elif document_type == "PDF":
        documents = parse_pdf(file_path)
    elif document_type == "PPTX":
        documents = parse_pptx(file_path)
    elif is_image:
        documents = parse_image(file_path)
        query_text = "Summarize the visual content previously extracted from this image."
    elif is_tabular:
        documents = parse_tabular_data(file_path)
        query_text = "Summarize the key insights and patterns from the provided tabular data."
    elif is_word:
        documents = parse_docx(file_path)
    else:
        print(f"Unsupported document type: {document_type}")
        return

    if documents is None or not documents:
        print(f"Failed to load or parse the {document_type} file. Skipping.")
        return

    print(f"{document_type} file parsed successfully and loaded into documents.")

    Settings.llm = LangChainLLM(llm=LangchainOllama(model=llm_model_to_use))
    Settings.embed_model = LangchainEmbedding(OllamaEmbeddings(model=OLLAMA_EMBED_MODEL))
    Settings.node_parser = SentenceSplitter(chunk_size=1024)
    print(f"LlamaIndex Settings configured for {document_type} with LLM: {llm_model_to_use}")

    if not os.path.exists(storage_dir) or not os.listdir(storage_dir):
        print(f"Building index and persisting to storage in: {storage_dir}")
        index = VectorStoreIndex.from_documents(documents)
        index.storage_context.persist(persist_dir=storage_dir)
        print(f"Index built and persisted to {storage_dir}")
    else:
        print(f"Loading index from existing storage: {storage_dir}")
        storage_context = StorageContext.from_defaults(persist_dir=storage_dir)
        index = load_index_from_storage(storage_context)
        print("Index loaded from storage.")

    print(f"\nQuerying the index for {document_type}: '{query_text}'")
    query_engine = index.as_query_engine(similarity_top_k=3)
    response = query_engine.query(query_text)

    print(f"\n=== Summary for {document_type} ===\n")
    print(response)
    print(f"\n--- Finished processing {document_type} ---")

def main():
    process_document(MSG_FILE_PATH, STORAGE_DIR_MSG, "Outlook MSG")
    process_document(PDF_FILE_PATH, STORAGE_DIR_PDF, "PDF")
    process_document(PPTX_FILE_PATH, STORAGE_DIR_PPTX, "PPTX")
    process_document(DOCX_FILE_PATH, STORAGE_DIR_DOCX, "Word Document", is_word=True)

    mime_type, _ = mimetypes.guess_type(IMAGE_FILE_PATH)
    if mime_type and mime_type.startswith("image"):
        process_document(IMAGE_FILE_PATH, STORAGE_DIR_IMG, "Image", is_image=True)

    process_document(TABULAR_FILE_PATH, STORAGE_DIR_TABULAR, "Tabular Data", is_tabular=True)


if __name__ == "__main__":
    main()